from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- palette
DARK = RGBColor(0x0b, 0x13, 0x20)
BLUE = RGBColor(0x25, 0x63, 0xeb)
BLUE_LIGHT = RGBColor(0x93, 0xc5, 0xfd)
WHITE = RGBColor(0xff, 0xff, 0xff)
TEXT = RGBColor(0x1a, 0x1a, 0x1a)
MUTED = RGBColor(0x94, 0xa3, 0xb8)
MUTED_DARK = RGBColor(0x47, 0x55, 0x69)
HEADER_BG = RGBColor(0xef, 0xf6, 0xff)
BORDER = RGBColor(0xcb, 0xd5, 0xe1)
CODE_BG = RGBColor(0x0f, 0x17, 0x2a)
CODE_TEXT = RGBColor(0xe2, 0xe8, 0xf0)

SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header(slide, title, dark=False, tag=None):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.32), Inches(12.1), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE if dark else TEXT
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.12), Inches(12.13), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    if tag:
        ftb = slide.shapes.add_textbox(Inches(11.6), Inches(7.05), Inches(1.6), Inches(0.35))
        fp = ftb.text_frame.paragraphs[0]
        fp.alignment = PP_ALIGN.RIGHT
        fr = fp.add_run()
        fr.text = tag
        fr.font.size = Pt(9)
        fr.font.color.rgb = MUTED if dark else MUTED_DARK
    return tb


def add_bullets(slide, items, left, top, width, height, font_size=15, color=TEXT, space_after=6):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = min(level, 8)
        bullet = "•" if level == 0 else "–"
        run = p.add_run()
        run.text = f"{bullet}  {text}"
        run.font.size = Pt(font_size - level * 1.5)
        run.font.color.rgb = color
        p.space_after = Pt(space_after)
    return tb


def add_table(slide, data, left, top, width, height, col_widths=None, font_size=12,
               header_align=None):
    rows, cols = len(data), len(data[0])
    shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            tf = cell.text_frame
            tf.word_wrap = True
            cell.margin_left = Inches(0.09)
            cell.margin_right = Inches(0.09)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = HEADER_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table


def add_pipeline(slide, items, left, top, box_w, box_h, gap):
    x = left
    for i, text in enumerate(items):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(box_w), Inches(box_h))
        box.fill.solid()
        box.fill.fore_color.rgb = HEADER_BG
        box.line.color.rgb = BLUE
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        for j, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = line
            run.font.size = Pt(13 if j == 0 else 10.5)
            run.font.bold = (j == 0)
            run.font.color.rgb = TEXT
        x += box_w
        if i < len(items) - 1:
            arrow = slide.shapes.add_textbox(Inches(x), Inches(top + box_h / 2 - 0.3), Inches(gap), Inches(0.6))
            ap = arrow.text_frame.paragraphs[0]
            ap.alignment = PP_ALIGN.CENTER
            ar = ap.add_run()
            ar.text = "→"
            ar.font.size = Pt(28)
            ar.font.bold = True
            ar.font.color.rgb = BLUE
            x += gap


def add_code_box(slide, text, left, top, width, height, font_size=14):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = CODE_TEXT
    run.font.name = "Menlo"


LOGO = "../../frontend/public/logo/floodguard-icon-256.png"
ARCH = "architecture.png"


# ================================================================== TITLE
s = add_slide()
set_bg(s, DARK)
s.shapes.add_picture(LOGO, Inches(5.92), Inches(0.75), height=Inches(1.7))

tb = s.shapes.add_textbox(Inches(1), Inches(2.7), Inches(11.33), Inches(1.2))
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "FloodGuard AI"
r.font.size = Pt(54)
r.font.bold = True
r.font.color.rgb = WHITE

tb = s.shapes.add_textbox(Inches(1), Inches(3.85), Inches(11.33), Inches(0.6))
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "From Initial-Round Model to a Deployed, AI-Powered MLOps Product"
r.font.size = Pt(20)
r.font.color.rgb = BLUE_LIGHT

tb = s.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(0.5))
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "ML Opsidian: Genesis — Final Round Submission"
r.font.size = Pt(15)
r.font.color.rgb = MUTED

tb = s.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.33), Inches(0.5))
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Team: [Add team name]   |   June 2026"
r.font.size = Pt(13)
r.font.color.rgb = MUTED


# ================================================================== 1. TEAM
s = add_slide()
add_header(s, "1. Team Introduction", tag="1 / 9")

add_bullets(s, [
    ('"We are [Team Name] — we turned our Initial Round flood-risk model into a '
     'deployed, AI-powered product (web + mobile) with full MLOps."', 0),
], 0.6, 1.45, 12.1, 0.9, font_size=18, color=MUTED_DARK)

add_table(s, [
    ["Role", "Member"],
    ["ML / Feature Engineering", "[Add name]"],
    ["MLOps / Deployment", "[Add name]"],
    ["Frontend / AI Advisor", "[Add name]"],
    ["Presentation", "[Add name]"],
], left=1.8, top=2.7, width=9.7, height=2.6, col_widths=[4.0, 5.7], font_size=16)

add_bullets(s, [
    ("Replace the placeholders above with your team's actual names and roles before presenting.", 0),
], 0.6, 5.7, 12.1, 0.6, font_size=12, color=MUTED_DARK)


# ================================================================== 2. PROBLEM UNDERSTANDING
s = add_slide()
add_header(s, "2. Problem Understanding", tag="2 / 9")

add_bullets(s, [
    ("Predict a continuous flood_risk_score in [0, 1] for locations across Sri Lanka's 25 districts", 0),
    ("Dataset: 20,886 training rows · 5,300 test rows · 47 columns", 0),
], 0.6, 1.35, 12.1, 1.0, font_size=15)

add_code_box(s, "score = (0.5 × MAE + 0.5 × RMSE) × (1 + max(0, 1 − explained_variance))",
             0.6, 2.35, 12.13, 0.6, font_size=15)

add_bullets(s, [
    ("Lower is better — rewards both accuracy AND calibration (matching the variance, not just the mean)", 0),
], 0.6, 3.0, 12.1, 0.5, font_size=13, color=MUTED_DARK)

add_table(s, [
    ["Category", "Example columns"],
    ["Geography", "latitude, longitude, district, elevation, distance to river"],
    ["Climate", "7-day & monthly rainfall, drainage index, seasonal index"],
    ["Land", "land cover, soil type, NDVI / NDWI, built-up %"],
    ["Infrastructure & Socioeconomics", "population density, infrastructure score, distance to hospital"],
    ["Event context", "flood occurrence, inundated area, historical flood count"],
], left=0.6, top=3.65, width=12.13, height=2.6, col_widths=[3.6, 8.53], font_size=13)

add_bullets(s, [
    ("Key insight: district dominates feature importance — district_te is the #1 feature by gain in "
     "every model version. The data is synthetic/noisy → ensembling + district-level aggregates win.", 0),
], 0.6, 6.45, 12.1, 0.8, font_size=13, color=MUTED_DARK)


# ================================================================== 3. INITIAL ROUND
s = add_slide()
add_header(s, "3. Initial Round Approach — solution_v10.py", tag="3 / 9")

add_pipeline(s, [
    "5-Model Ensemble\nLightGBM · LightGBM-DART\nCatBoost · XGBoost\nExtraTrees (Optuna-tuned)",
    "Ridge Meta-Learner\nStacks 10-fold × 3-seed\nOOF predictions",
    "Calibration Shift\nShift mean to match\ntrain target mean",
    "Custom Metric\n≈ 0.383",
], left=0.6, top=1.5, box_w=2.65, box_h=1.9, gap=0.55)

add_bullets(s, [
    ("~150 engineered features: district stats/quantiles/ranks, KMeans clusters (k = 5, 10, 20), "
     "haversine-KNN target stats, target encoding, interaction features", 0),
    ("10-fold × 3-seed cross-validation", 0),
    ("10 rounds of pseudo-labeling on the test set to squeeze out additional signal", 0),
    ("Post-processing: shift the final prediction mean to exactly match the training target mean "
     "(the calibration insight that carries through to the Final Round)", 0),
], 0.6, 3.8, 12.1, 3.0, font_size=15, space_after=10)


# ================================================================== 4. IMPROVEMENTS MADE
s = add_slide()
add_header(s, "4. Improvements Made: v10 → FloodGuard AI", tag="4 / 9")

add_table(s, [
    ["v10 (batch)", "FloodGuard AI (FeatureEngineer)"],
    ["District stats via groupby over all rows", "Precomputed per-district lookup table, joined onto new rows"],
    ["KMeans fit and .predict() on the same data", "KMeans models fit once, saved, .predict() on the new row"],
    ["KNN target stats via BallTree over all rows", "Same BallTree, built once at fit time, queried per new row"],
    ["Target encoding via K-fold OOF means", "OOF means for training; full-train means for new rows"],
    ["*_qmap columns (no raw counterpart)", "Empirical monotonic map learned at fit time, np.interp at transform time"],
    ["Composite indices (seasonal_index, etc.) assumed present", "Default to the district mean (or global mean) if the user omits them"],
], left=0.6, top=1.35, width=12.13, height=2.85, col_widths=[5.6, 6.53], font_size=11.5)

add_bullets(s, [
    ("Reduced to a 3-model ensemble (LightGBM + CatBoost + XGBoost), 5-fold CV, fixed hyperparameters, "
     "no pseudo-labeling → ~45-second, fully reproducible training run — 223 features, 14 MB of artifacts", 0),
    ("Kept the two highest-leverage ideas: Ridge meta-stacking (weights: lgb=0.229, cat=0.713, xgb=0.134) "
     "and mean-calibration shift (shift ≈ 0 — stacked mean already matches the train mean of 0.4780)", 0),
    ("Result: OOF custom metric 0.4071 (vs 0.383 for the heavier Initial Round search) — a documented "
     "trade-off of leaderboard score for deployability", 0),
    ("NEW: an AI Risk Advisor (Claude claude-haiku-4-5, with a deterministic template fallback) turns the "
     "score + top contributing features into a plain-language report and recommendations", 0),
], 0.6, 4.35, 12.1, 2.9, font_size=13, space_after=8)


# ================================================================== 5. SYSTEM ARCHITECTURE
s = add_slide()
add_header(s, "5. System Architecture", tag="5 / 9")
s.shapes.add_picture(ARCH, Inches(2.55), Inches(1.25), height=Inches(5.1))
add_bullets(s, [
    ("One Docker image, one Railway URL serves the REST API, the prediction form, and the monitoring "
     "dashboard — judges don't need to set up anything", 0),
], 9.55, 1.4, 3.5, 1.6, font_size=12)
add_bullets(s, [
    ("A companion Expo / React Native mobile app (Android APK + Expo Go) talks to the same live backend "
     "and adds a geolocation-based Live Risk Preview on both clients", 0),
], 9.55, 3.1, 3.5, 1.8, font_size=12)
add_bullets(s, [
    ("Offline path: src/train.py → MLflow (mlruns/) → models/v1/ artifacts loaded at API startup", 0),
], 9.55, 4.95, 3.5, 1.6, font_size=12)


# ================================================================== 6. MLOPS WORKFLOW
s = add_slide()
add_header(s, "6. MLOps Workflow", tag="6 / 9")

add_table(s, [
    ["Pillar", "What We Built"],
    ["Data Pipeline", "src/data_validation.py runs schema/null/range checks before training, plus per-request out-of-distribution flags"],
    ["Model Management", "MLflow experiment tracking (mlruns/) logs every run's hyperparameters and metrics; models/v1/metadata.json records feature count, CV metrics, Ridge weights — versioned for v2, v3, …"],
    ["Deployment", "Dockerfile → Railway (railway up), /health readiness check, optional ANTHROPIC_API_KEY; Render documented as an alternative; Android APK via EAS Build"],
    ["Monitoring", "Every prediction logged to SQLite (score, category, latency, OOD flags); /monitoring/stats + /dashboard (Chart.js) auto-refresh every 15s"],
    ["CI/CD (bonus)", "GitHub Actions runs the full pytest suite (feature pipeline, inference, API) and builds the Docker image on every push/PR"],
], left=0.6, top=1.35, width=12.13, height=4.5, col_widths=[2.6, 9.53], font_size=13)


# ================================================================== 7. DEMONSTRATION
s = add_slide()
add_header(s, "7. Demonstration (Live + Backup Recording)", tag="7 / 9")

add_bullets(s, [
    ("Landing page (/) — hero, \"how it works\", live model snapshot from /model/info; the geolocation-based "
     "Live Risk Preview shows the risk for the user's own district instantly, no form needed", 0),
    ("Prediction flow — fill the form for a flood-prone location → risk gauge, risk category, model "
     "breakdown (lgb / cat / xgb), top contributing factors, AI-generated risk report with recommendations", 0),
    ("Edge case — submit an out-of-range value (e.g. negative rainfall) → the OOD flag appears in the response", 0),
    ("Feedback — rate the prediction with the star widget → POST /feedback → SQLite", 0),
    ("Dashboard (/dashboard) — prediction volume, score histogram, category/district breakdowns, latency, "
     "and the feedback just submitted, all live", 0),
    ("Mobile app (optional) — the same flow on the Android APK / Expo Go, including the geolocation "
     "Live Risk Preview on the Home screen", 0),
], 0.6, 1.4, 12.1, 4.6, font_size=15, space_after=12)

add_bullets(s, [
    ("Rehearse this end-to-end against the live Railway URL, not just localhost — and have a screen-recording "
     "ready as a fallback if the live demo or network fails.", 0),
], 0.6, 6.3, 12.1, 0.8, font_size=12, color=MUTED_DARK)


# ================================================================== 8. CHALLENGES
s = add_slide()
add_header(s, "8. Challenges Faced", tag="8 / 9")

add_bullets(s, [
    ("Making a batch feature pipeline servable for one record — v10's district stats, KMeans, KNN, and "
     "target encodings were all computed over the entire training set at once; redesigning each as a "
     "fit-once / apply-to-one-row component (FeatureEngineer) was the single biggest engineering effort "
     "of the Final Round.", 0),
    ("Trading leaderboard score for deployability — dropping pseudo-labeling, 2 of 5 base models, and "
     "going from 10-fold × 3-seed to 5-fold cost ~0.024 on the custom metric (0.383 → 0.407) — a "
     "deliberate, documented trade-off rather than an oversight.", 0),
    ("Handling \"expected\" data noise without over-engineering — negative rainfall/distance values produce "
     "NaN after log1p(); rather than imputing or dropping rows, we leaned on LightGBM/CatBoost/XGBoost's "
     "native NaN handling and added flag_out_of_distribution() so these cases are visible in production "
     "rather than silently handled.", 0),
], 0.6, 1.5, 12.1, 5.0, font_size=16, space_after=22)


# ================================================================== 9. FUTURE IMPROVEMENTS
s = add_slide()
add_header(s, "9. Future Improvements", tag="9 / 9")

add_bullets(s, [
    ("Periodic retraining as new labeled flood events become available, versioned as models/v2/, v3/, "
     "… and compared via the MLflow run history", 0),
    ("Lightweight, request-time-feasible pseudo-labeling using accumulated predictions + user feedback "
     "as a weak-label source", 0),
    ("A model-drift alert (rolling flag_out_of_distribution rate vs. a baseline) surfaced on the dashboard", 0),
    ("Expanding the AI Risk Advisor to support Sinhala / Tamil for end users in the regions the model covers", 0),
    ("Publishing the mobile app to the Play Store / TestFlight, plus background location + push "
     "notifications when a user's district risk crosses a threshold", 0),
], 0.6, 1.45, 12.1, 3.6, font_size=15, space_after=10)

add_code_box(s,
    '"FloodGuard AI shows the Initial Round model didn’t just survive the transition to production — '
    'its core ideas (ensembling, stacking, calibration) are exactly what made that transition work."',
    0.6, 5.6, 12.13, 1.3, font_size=15)


# ================================================================== THANK YOU
s = add_slide()
set_bg(s, DARK)
s.shapes.add_picture(LOGO, Inches(5.92), Inches(0.7), height=Inches(1.5))

tb = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.0))
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Thank You — Questions?"
r.font.size = Pt(40)
r.font.bold = True
r.font.color.rgb = WHITE

links = [
    ("GitHub:", "github.com/malika1234m/flood-guard-ai"),
    ("Live Demo:", "floodguard-production.up.railway.app"),
    ("Mobile:", "Android APK (EAS internal build) — link in README"),
]
tb = s.shapes.add_textbox(Inches(2.5), Inches(4.0), Inches(8.33), Inches(2.0))
tf = tb.text_frame
tf.word_wrap = True
for i, (label, val) in enumerate(links):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = f"{label} "
    r1.font.size = Pt(16)
    r1.font.bold = True
    r1.font.color.rgb = BLUE_LIGHT
    r2 = p.add_run()
    r2.text = val
    r2.font.size = Pt(16)
    r2.font.color.rgb = WHITE
    p.space_after = Pt(10)


OUT = "../presentation.pptx"
prs.save(OUT)
print("wrote", OUT, "-", len(prs.slides._sldIdLst), "slides")
